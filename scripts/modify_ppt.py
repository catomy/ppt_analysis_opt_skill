#!/usr/bin/env python3
"""
Modify PowerPoint presentation content.

This script applies changes to PPTX files based on structured modification data.
"""

import sys
import json
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def apply_slide_changes(slide, target_type, changes):
    """Apply changes to a single slide by paragraph index."""
    # Sort changes by paragraph_index (descending) to avoid index shifting
    changes_sorted = sorted(changes, key=lambda c: c.get('paragraph_index', 0), reverse=True)

    for change in changes_sorted:
        change_type = change.get('type')
        paragraph_index = change.get('paragraph_index', 0)
        new_text = change.get('new_text', '')
        old_text = change.get('old_text', '')
        match_mode = change.get('match_mode', 'exact')

        if change_type == 'replace_by_index':
            # Find the target shape by position
            if target_type == 'title':
                ordered_shapes = get_shapes_in_order(slide)
                title_shape = identify_title_shape(ordered_shapes)
                if title_shape and title_shape.has_text_frame:
                    replaced = replace_paragraph_by_index(
                        title_shape.text_frame,
                        0,
                        new_text,
                        expected_old_text=old_text,
                        match_mode=match_mode,
                    )
                    if not replaced and old_text:
                        find_and_replace_by_old_text(slide, old_text, new_text)

            elif target_type == 'content':
                ordered_shapes = get_shapes_in_order(slide)
                title_shape = identify_title_shape(ordered_shapes)
                content_shapes = [s for s in ordered_shapes if s != title_shape and getattr(s, 'has_text_frame', False) and s.text.strip()]

                # Find the paragraph by counting
                current_para_index = 0
                for shape in content_shapes:
                    # Count paragraphs in this shape
                    para_count = len([p for p in shape.text_frame.paragraphs if p.text.strip()])

                    if current_para_index + para_count > paragraph_index:
                        # The target paragraph is in this shape
                        local_index = paragraph_index - current_para_index
                        replaced = replace_paragraph_by_index(
                            shape.text_frame,
                            local_index,
                            new_text,
                            expected_old_text=old_text,
                            match_mode=match_mode,
                        )
                        if not replaced and old_text:
                            find_and_replace_by_old_text(slide, old_text, new_text)
                        break
                    else:
                        current_para_index += para_count

        elif change_type == 'replace_by_shape_paragraph':
            shape_index = change.get('shape_index', None)
            paragraph_index_in_shape = change.get('paragraph_index_in_shape', None)
            nonempty_index_in_shape = change.get('nonempty_index_in_shape', None)

            ordered_shapes = get_shapes_in_order(slide)
            if not isinstance(shape_index, int) or not (0 <= shape_index < len(ordered_shapes)):
                if old_text:
                    find_and_replace_by_old_text(slide, old_text, new_text)
                continue

            shape = ordered_shapes[shape_index]
            if not getattr(shape, 'has_text_frame', False):
                if old_text:
                    find_and_replace_by_old_text(slide, old_text, new_text)
                continue

            replaced = False
            if isinstance(paragraph_index_in_shape, int) and 0 <= paragraph_index_in_shape < len(shape.text_frame.paragraphs):
                target_paragraph = shape.text_frame.paragraphs[paragraph_index_in_shape]
                replaced = replace_paragraph(target_paragraph, new_text, expected_old_text=old_text, match_mode=match_mode)
            else:
                if isinstance(nonempty_index_in_shape, int):
                    replaced = replace_paragraph_by_index(
                        shape.text_frame,
                        nonempty_index_in_shape,
                        new_text,
                        expected_old_text=old_text,
                        match_mode=match_mode,
                    )

            if not replaced and old_text:
                find_and_replace_by_old_text(slide, old_text, new_text, preferred_shape_index=shape_index)

        elif change_type == 'replace_text':
            # Legacy support: text-based replacement
            old_text = change['old_text']
            replaced = False

            for shape in slide.shapes:
                if hasattr(shape, 'text') and old_text in shape.text and new_text not in shape.text:
                    shape.text = shape.text.replace(old_text, new_text)
                    replaced = True
                    break
                if replaced:
                    break

        elif change_type == 'add_textbox':
            # Add a new text box
            left = Inches(change.get('left', 1))
            top = Inches(change.get('top', 1))
            width = Inches(change.get('width', 5))
            height = Inches(change.get('height', 1))

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Add text
            p = text_frame.paragraphs[0]
            p.text = change['text']
            p.font.size = Pt(change.get('font_size', 18))

        elif change_type == 'delete_shape':
            # Delete shape containing specific text
            for shape in list(slide.shapes):
                if hasattr(shape, 'text') and change['text'] in shape.text:
                    sp = shape.element
                    sp.getparent().remove(sp)


def get_shapes_in_order(slide):
    shapes = []
    for shape in slide.shapes:
        shapes.extend(flatten_shape(shape))
    shapes.sort(key=lambda s: (s.top, s.left))
    return shapes


def flatten_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        flattened = []
        for child in shape.shapes:
            flattened.extend(flatten_shape(child))
        return flattened

    if getattr(shape, 'has_table', False):
        return [shape]

    if getattr(shape, 'has_text_frame', False) and shape.text.strip():
        return [shape]

    return []


def identify_title_shape(ordered_shapes):
    top_shapes = []
    for shape in ordered_shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        if shape.top >= 350000:
            continue
        if not shape.text.strip():
            continue

        first_para = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else None
        if not first_para:
            continue
        if not first_para.runs:
            continue
        font_size = getattr(first_para.runs[0].font.size, 'pt', 0) if first_para.runs[0].font.size else 0
        if font_size > 20:
            top_shapes.append((shape, font_size))

    if top_shapes:
        return max(top_shapes, key=lambda x: x[1])[0]

    for shape in ordered_shapes:
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                if shape.placeholder_format.type in [1, 14]:
                    return shape
        except ValueError:
            continue

    for shape in ordered_shapes:
        if getattr(shape, 'has_text_frame', False) and shape.text.strip():
            return shape
    return None


def normalize_text(text: str) -> str:
    return (text or '').strip()


def replace_paragraph(paragraph, new_text: str, expected_old_text: str = '', match_mode: str = 'exact') -> bool:
    current_text = normalize_text(getattr(paragraph, 'text', ''))
    expected = normalize_text(expected_old_text)
    if expected:
        if match_mode == 'contains':
            if expected not in current_text:
                return False
        else:
            if current_text != expected:
                return False

    if paragraph.runs:
        paragraph.runs[0].text = new_text
        for run in list(paragraph.runs)[1:]:
            r = run.element
            r.getparent().remove(r)
    else:
        paragraph.text = new_text
    return True


def replace_paragraph_by_index(text_frame, paragraph_index, new_text, expected_old_text: str = '', match_mode: str = 'exact'):
    """Replace a paragraph's text by index while preserving formatting."""
    # Get non-empty paragraphs
    paragraphs = [p for p in text_frame.paragraphs if p.text.strip()]

    if 0 <= paragraph_index < len(paragraphs):
        para = paragraphs[paragraph_index]
        return replace_paragraph(para, new_text, expected_old_text=expected_old_text, match_mode=match_mode)
    return False


def find_and_replace_by_old_text(slide, old_text: str, new_text: str, preferred_shape_index: int | None = None) -> bool:
    expected = normalize_text(old_text)
    if not expected:
        return False

    ordered_shapes = get_shapes_in_order(slide)
    if isinstance(preferred_shape_index, int) and 0 <= preferred_shape_index < len(ordered_shapes):
        preferred_shape = ordered_shapes[preferred_shape_index]
        if getattr(preferred_shape, 'has_text_frame', False):
            for paragraph in preferred_shape.text_frame.paragraphs:
                current_text = normalize_text(getattr(paragraph, 'text', ''))
                if expected == current_text and normalize_text(new_text) != current_text:
                    return replace_paragraph(paragraph, new_text, expected_old_text=old_text, match_mode='exact')

    for shape in ordered_shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            current_text = normalize_text(getattr(paragraph, 'text', ''))
            if expected == current_text and normalize_text(new_text) != current_text:
                return replace_paragraph(paragraph, new_text, expected_old_text=old_text, match_mode='exact')
    return False


def apply_modifications(pptx_path, modifications):
    """
    Apply structured modifications to a presentation.

    Returns:
        Presentation object on success, or dict with 'error' key on failure.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {'error': f'Failed to open presentation: {str(e)}'}

    style_settings = extract_style_settings(modifications)

    # Apply slide-specific changes
    for mod in modifications:
        slide_index = mod.get('slide_index')
        target_type = mod.get('target_type', 'content')  # 'title' or 'content'
        # slide_index 是 1-based，转换为 0-based 访问
        array_index = slide_index - 1 if slide_index is not None else None
        if array_index is not None and 0 <= array_index < len(prs.slides):
            slide = prs.slides[array_index]
            changes = mod.get('changes', [])
            apply_slide_changes(slide, target_type, changes)

    normalize_presentation_style(prs, style_settings)
    return prs


def extract_style_settings(modifications):
    settings = {
        'enabled': True,
        'default_font_name': '微软雅黑',
        'title_font_size': 28,
        'min_title_font_size': 24,
        'max_title_font_size': 30,
        'unify_title_color': True,
        'title_color_rgb': None,
        'enforce_title_single_line': True,
    }

    for mod in modifications or []:
        if mod.get('target_type') != 'style':
            continue
        for change in mod.get('changes', []):
            if change.get('type') != 'normalize_style':
                continue
            settings.update({k: v for k, v in change.items() if k != 'type'})
            rgb = change.get('title_color_rgb', None)
            if isinstance(rgb, (list, tuple)) and len(rgb) == 3 and all(isinstance(x, int) for x in rgb):
                settings['title_color_rgb'] = tuple(rgb)
    return settings


def normalize_presentation_style(prs, settings):
    if not settings or not settings.get('enabled', True):
        return

    default_font_name = settings.get('default_font_name', '微软雅黑')
    title_font_size = settings.get('title_font_size', 28)
    min_title_font_size = settings.get('min_title_font_size', 24)
    max_title_font_size = settings.get('max_title_font_size', 30)
    unify_title_color = settings.get('unify_title_color', True)
    enforce_title_single_line = settings.get('enforce_title_single_line', True)

    title_font_size = max(min_title_font_size, min(max_title_font_size, title_font_size))

    title_color_rgb = settings.get('title_color_rgb', None)
    if unify_title_color and title_color_rgb is None:
        title_color_rgb = detect_canonical_title_rgb(prs)

    for slide in prs.slides:
        ordered_shapes = get_shapes_in_order(slide)
        title_shape = identify_title_shape(ordered_shapes)

        for shape in slide.shapes:
            apply_default_font_to_shape(shape, default_font_name)

        if title_shape and getattr(title_shape, 'has_text_frame', False):
            if enforce_title_single_line:
                normalize_title_text_single_line(title_shape)
            apply_title_style_to_shape(title_shape, default_font_name, title_font_size, title_color_rgb)


def apply_default_font_to_shape(shape, default_font_name):
    if getattr(shape, 'has_table', False):
        for row in shape.table.rows:
            for cell in row.cells:
                apply_default_font_to_text_frame(cell.text_frame, default_font_name)
        return

    if getattr(shape, 'has_text_frame', False):
        apply_default_font_to_text_frame(shape.text_frame, default_font_name)


def apply_default_font_to_text_frame(text_frame, default_font_name):
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            for run in paragraph.runs:
                run.font.name = default_font_name
        else:
            paragraph.font.name = default_font_name


def apply_title_style_to_shape(shape, default_font_name, title_font_size, title_color_rgb):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = default_font_name
            run.font.size = Pt(title_font_size)
            run.font.bold = True
            if title_color_rgb is not None:
                run.font.color.rgb = RGBColor(*title_color_rgb)


def normalize_title_text_single_line(shape):
    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.text:
            continue
        normalized = normalize_text(paragraph.text).replace('\n', '').replace('\r', '')
        if paragraph.runs:
            paragraph.runs[0].text = normalized
            for run in list(paragraph.runs)[1:]:
                r = run.element
                r.getparent().remove(r)
        else:
            paragraph.text = normalized


def detect_canonical_title_rgb(prs):
    for slide in prs.slides:
        ordered_shapes = get_shapes_in_order(slide)
        title_shape = identify_title_shape(ordered_shapes)
        if not title_shape or not getattr(title_shape, 'has_text_frame', False):
            continue
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                color = getattr(run.font, 'color', None)
                rgb = getattr(color, 'rgb', None) if color else None
                if rgb is not None:
                    return (rgb[0], rgb[1], rgb[2])
    return None


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(json.dumps({
            'error': 'Usage: modify_ppt.py <pptx_file_path> <modifications_json> <output_path>'
        }))
        sys.exit(1)

    pptx_path = sys.argv[1]
    modifications_json = sys.argv[2]
    output_path = sys.argv[3] if len(sys.argv) > 3 else pptx_path

    try:
        # Try to load from file first
        try:
            with open(modifications_json, 'r', encoding='utf-8') as f:
                modifications = json.load(f)
        except FileNotFoundError:
            # If not a file, try to parse as JSON string
            modifications = json.loads(modifications_json)

        prs = apply_modifications(pptx_path, modifications)

        # Check for errors
        if isinstance(prs, dict):
            print(json.dumps(prs, ensure_ascii=False))
            sys.exit(1)

        # Save modified presentation (prs is Presentation object here)
        prs.save(output_path)
        print(json.dumps({
            'success': True,
            'output_file': output_path,
            'message': 'Presentation modified successfully'
        }, ensure_ascii=False, indent=2))

    except json.JSONDecodeError as e:
        print(json.dumps({'error': f'Invalid JSON: {str(e)}'}, ensure_ascii=False))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({'error': f'Unexpected error: {str(e)}'}, ensure_ascii=False))
        sys.exit(1)
