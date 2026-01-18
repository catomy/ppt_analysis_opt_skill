import json
import os
import tempfile
import unittest

from pptx import Presentation
from pptx.dml.color import RGBColor

from utils.text_extraction import PPTXTextExtractor
from scripts.modify_ppt import apply_modifications


def iter_all_paragraph_texts(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = (paragraph.text or "").strip()
                    if text:
                        yield text


class TestPptLogicAnalyzerRoundTrip(unittest.TestCase):
    def test_extractor_to_json_is_serializable(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, "input.pptx")

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "旧标题"
            body = slide.placeholders[1].text_frame
            body.clear()
            body.paragraphs[0].text = "第一条"
            body.add_paragraph().text = "第二条"
            prs.save(input_path)

            extractor = PPTXTextExtractor(input_path)
            payload = extractor.to_json()
            data = json.loads(payload)

            self.assertEqual(data["file"], input_path)
            self.assertEqual(data["total_slides"], 1)
            self.assertIn("content", data["slides"][0])

    def test_modify_replace_by_shape_paragraph(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, "input.pptx")
            output_path = os.path.join(temp_dir, "output.pptx")

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "旧标题"

            body = slide.placeholders[1].text_frame
            body.clear()
            body.paragraphs[0].text = "第一条"
            body.add_paragraph().text = "第二条"

            slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height).text_frame.text = "第三条"
            prs.save(input_path)

            extracted = PPTXTextExtractor(input_path).extract_all_slides()
            first_slide = extracted["slides"][0]
            target_entry = next(p for p in first_slide["content"] if p["text"] == "第二条")

            modifications = [
                {
                    "slide_index": 1,
                    "target_type": "content",
                    "changes": [
                        {
                            "type": "replace_by_shape_paragraph",
                            "shape_index": target_entry["shape_index"],
                            "paragraph_index_in_shape": target_entry["paragraph_index_in_shape"],
                            "nonempty_index_in_shape": target_entry["nonempty_index_in_shape"],
                            "paragraph_index": target_entry["paragraph_index"],
                            "old_text": target_entry["text"],
                            "new_text": "第二条（已替换）",
                        }
                    ],
                }
            ]

            updated_prs = apply_modifications(input_path, modifications)
            self.assertTrue(hasattr(updated_prs, "slides"))
            self.assertTrue(hasattr(updated_prs, "save"))
            updated_prs.save(output_path)

            updated = Presentation(output_path)
            all_texts = list(iter_all_paragraph_texts(updated))
            self.assertIn("第二条（已替换）", all_texts)

    def test_modify_replace_by_index_with_old_text_fallback(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, "input.pptx")

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "旧标题"

            body = slide.placeholders[1].text_frame
            body.clear()
            body.paragraphs[0].text = "第一条"
            body.add_paragraph().text = "第二条"
            prs.save(input_path)

            modifications = [
                {
                    "slide_index": 1,
                    "target_type": "content",
                    "changes": [
                        {
                            "type": "replace_by_index",
                            "paragraph_index": 0,
                            "old_text": "第二条",
                            "new_text": "第二条（fallback替换）",
                        }
                    ],
                }
            ]

            updated_prs = apply_modifications(input_path, modifications)
            all_texts = list(iter_all_paragraph_texts(updated_prs))
            self.assertIn("第二条（fallback替换）", all_texts)

    def test_style_normalization_default_font_and_title_style(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, "input.pptx")

            prs = Presentation()

            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "第一张标题"
            title_run = slide1.shapes.title.text_frame.paragraphs[0].runs[0]
            title_run.font.color.rgb = RGBColor(10, 20, 30)

            body1 = slide1.placeholders[1].text_frame
            body1.clear()
            body1.paragraphs[0].text = "正文A"

            slide2 = prs.slides.add_slide(prs.slide_layouts[0])
            slide2.shapes.title.text = "第二张标题"
            body2 = slide2.placeholders[1].text_frame
            body2.clear()
            body2.paragraphs[0].text = "正文B"

            prs.save(input_path)

            updated_prs = apply_modifications(input_path, [])

            for slide in updated_prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                self.assertEqual(run.font.name, "微软雅黑")

            slide1_title = updated_prs.slides[0].shapes.title
            slide2_title = updated_prs.slides[1].shapes.title

            slide1_title_run = slide1_title.text_frame.paragraphs[0].runs[0]
            slide2_title_run = slide2_title.text_frame.paragraphs[0].runs[0]

            self.assertEqual(slide1_title_run.font.size.pt, 28)
            self.assertEqual(slide2_title_run.font.size.pt, 28)
            self.assertTrue(slide1_title_run.font.bold)
            self.assertTrue(slide2_title_run.font.bold)
            self.assertEqual(slide1_title_run.font.color.rgb, slide2_title_run.font.color.rgb)


if __name__ == "__main__":
    unittest.main()

