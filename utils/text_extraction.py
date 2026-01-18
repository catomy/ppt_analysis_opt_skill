from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Optional
import json


class PPTXTextExtractor:
    """段落级PPT文本提取器"""

    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    def extract_all_slides(self) -> Dict[str, Any]:
        """提取所有幻灯片的结构化文本"""
        slides_data = []

        for idx, slide in enumerate(self.prs.slides):
            slide_data = self._extract_slide(slide, idx)
            slides_data.append(slide_data)

        return {
            'file': self.pptx_path,
            'total_slides': len(slides_data),
            'slides': slides_data
        }

    def _extract_slide(self, slide, slide_idx: int) -> Dict[str, Any]:
        """提取单页幻灯片内容"""
        # 按阅读顺序排序形状
        ordered_shapes = self._get_shapes_in_order(slide)

        # 分离标题和内容
        title_shape = self._identify_title(ordered_shapes)
        content_shapes = [s for s in ordered_shapes if s != title_shape]

        # 提取段落级内容
        slide_data = {
            'slide_number': slide_idx + 1,
            'slide_id': slide.slide_id,
            'title': self._extract_paragraphs_from_shape(title_shape) if title_shape else None,
            'title_meta': self._get_shape_meta(ordered_shapes, title_shape) if title_shape else None,
            'content': [],
            'content_shapes': [],
            'tables': [],
            'notes': self._extract_notes(slide),
            'has_header_footer': self._detect_header_footer(ordered_shapes)
        }

        global_paragraph_index = 0
        for shape in content_shapes:
            shape_meta = self._get_shape_meta(ordered_shapes, shape)
            shape_data = {
                'type': self._get_shape_type(shape),
                'position': {'top': shape.top, 'left': shape.left},
                'paragraphs': self._extract_paragraphs_from_shape(shape)
            }

            # 表格特殊处理
            if shape.has_table:
                table_data = self._extract_table_data(shape)
                shape_data['table_data'] = table_data
                slide_data['tables'].append({
                    **shape_meta,
                    **table_data
                })
            else:
                nonempty_index_in_shape = 0
                for paragraph in shape_data['paragraphs']:
                    slide_data['content'].append({
                        'paragraph_index': global_paragraph_index,
                        **shape_meta,
                        'paragraph_index_in_shape': paragraph.get('index'),
                        'nonempty_index_in_shape': nonempty_index_in_shape,
                        'text': paragraph.get('text', ''),
                        'level': paragraph.get('level'),
                        'alignment': paragraph.get('alignment'),
                        'runs': paragraph.get('runs', [])
                    })
                    global_paragraph_index += 1
                    nonempty_index_in_shape += 1

            slide_data['content_shapes'].append(shape_data)

        return slide_data

    def _extract_paragraphs_from_shape(self, shape) -> List[Dict]:
        """从形状中提取段落（含层级和格式）"""
        if not shape or not shape.has_text_frame:
            return []

        paragraphs = []
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            if not para.text.strip():
                continue

            # 提取run级格式
            runs_data = []
            for run in para.runs:
                if run.text.strip():
                    runs_data.append({
                        'text': run.text,
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'font_size': run.font.size.pt if run.font.size else None,
                        'font_name': run.font.name
                    })

            paragraphs.append({
                'index': para_idx,
                'text': para.text.strip(),
                'level': para.level,  # 0-8，层级
                'alignment': str(para.alignment) if para.alignment else None,
                'runs': runs_data
            })

        return paragraphs

    def _extract_table_data(self, shape) -> Dict:
        """提取表格数据"""
        table = shape.table
        cells = []

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if not cell.is_spanned:
                    cells.append({
                        'row': row_idx,
                        'col': col_idx,
                        'text': cell.text.strip()
                    })

        return {
            'rows': len(table.rows),
            'cols': len(table.columns),
            'cells': cells
        }

    def _get_shapes_in_order(self, slide) -> List:
        """按阅读顺序获取形状（文本、表格、组合形状子节点）"""
        shapes: List[Any] = []
        for shape in slide.shapes:
            shapes.extend(self._flatten_shape(shape))

        shapes.sort(key=lambda s: (s.top, s.left))
        return shapes

    def _flatten_shape(self, shape) -> List:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            flattened: List[Any] = []
            for child in shape.shapes:
                flattened.extend(self._flatten_shape(child))
            return flattened

        if getattr(shape, 'has_table', False):
            return [shape]

        if getattr(shape, 'has_text_frame', False) and shape.text.strip():
            return [shape]

        return []

    def _get_shape_meta(self, ordered_shapes: List, shape) -> Dict[str, Any]:
        if shape is None:
            return {'shape_index': None, 'shape_id': None}
        try:
            shape_index = ordered_shapes.index(shape)
        except ValueError:
            shape_index = None
        shape_id = getattr(shape, 'shape_id', None)
        return {'shape_index': shape_index, 'shape_id': shape_id}

    def _identify_title(self, shapes: List) -> Any:
        """识别标题形状（优化版）"""
        if not shapes:
            return None

        # 方法1: 位置判断（顶部，字体较大）
        top_shapes = []
        for shape in shapes:
            if shape.top < 350000:  # EMUs，扩大范围
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    first_para = shape.text_frame.paragraphs[0]
                    if first_para.runs:
                        font_size = getattr(first_para.runs[0].font.size, 'pt', 0) if first_para.runs[0].font.size else 0
                        if font_size > 20:  # 降低字体阈值
                            top_shapes.append((shape, font_size))

        # 返回字体最大的顶部形状
        if top_shapes:
            return max(top_shapes, key=lambda x: x[1])[0]

        # 方法2: 占位符判断
        for shape in shapes:
            try:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.type in [1, 14]:  # 标题占位符
                        return shape
            except ValueError:
                continue

        # 方法3: 第一个非空文本形状
        for shape in shapes:
            if shape.has_text_frame and shape.text.strip():
                return shape

        return None

    def _extract_notes(self, slide) -> str:
        """提取演讲者备注"""
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text.strip()
        return ""

    def _detect_header_footer(self, shapes: List) -> bool:
        """检测页眉页脚"""
        # 简化版：检查明显的页眉页脚位置
        slide_height = getattr(self.prs, 'slide_height', 6858000)

        for shape in shapes:
            # 底部10%或顶部10%且高度小
            if shape.top > slide_height * 0.9 or (shape.top < slide_height * 0.1 and shape.height < slide_height * 0.1):
                return True
        return False

    def _get_shape_type(self, shape) -> str:
        """获取形状类型"""
        if shape.has_table:
            return "table"
        elif shape.has_chart:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        else:
            try:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    return "placeholder"
            except ValueError:
                pass
            return "textbox"

    def to_json(self) -> str:
        """导出为JSON"""
        return json.dumps(self.extract_all_slides(), ensure_ascii=False, indent=2)


# CLI入口
if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print(json.dumps({'error': 'Usage: text_extraction.py <pptx_file>'}, ensure_ascii=False))
        sys.exit(1)

    extractor = PPTXTextExtractor(sys.argv[1])
    print(extractor.to_json())
